import fitz  # PyMuPDF
import docx
import io
import uuid
import os
import random
import base64
from PIL import Image
from groq import Groq
from fastapi import UploadFile
from langchain_text_splitters import RecursiveCharacterTextSplitter
import openpyxl
from pptx import Presentation
import json

from services.embedding_service import get_embeddings
from services.vector_service import store_chunks
from services.ocr_service import extract_text_with_vision_api
from services.image_preprocessing import preprocess_for_ocr

def get_groq_client():
    keys = os.environ.get("GROQ_API_KEYS", "")
    key_list = [k.strip() for k in keys.split(",") if k.strip()]
    if not key_list:
        raise ValueError("GROQ_API_KEYS not configured")
    api_key = random.choice(key_list)
    return Groq(api_key=api_key)

async def extract_text(file: UploadFile) -> str:
    """
    Multimodal Document & Image Extraction Pipeline.
    Supports PDF, DOCX, XLSX, PPTX, TXT, MD, CSV, PNG, JPG, JPEG, WEBP.
    """
    content = await file.read()
    filename = file.filename
    filename_lower = filename.lower()
    text = f"--- UPLOADED FILE: {filename} ---\n\n"
    
    try:
        if filename_lower.endswith(".pdf"):
            pdf_document = fitz.open(stream=content, filetype="pdf")
            total_pages = len(pdf_document)
            text += f"[DOCUMENT TYPE: PDF Document - {total_pages} Pages]\n\n"
            
            for page_num in range(total_pages):
                page = pdf_document.load_page(page_num)
                page_text = page.get_text()
                
                # Scanned page / no text -> Use Vision OCR
                if not page_text.strip():
                    pix = page.get_pixmap()
                    img_bytes = pix.tobytes("jpeg")
                    try:
                        page_text = extract_text_with_vision_api(img_bytes, mime_type="image/jpeg")
                    except Exception as e:
                        print(f"Vision OCR failed on PDF page {page_num+1}: {e}")
                        page_text = f"Page {page_num+1}: [Scanned page visual content]"
                
                text += f"=== Page {page_num+1} ===\n{page_text}\n\n"
                
        elif filename_lower.endswith(".docx"):
            doc = docx.Document(io.BytesIO(content))
            text += "[DOCUMENT TYPE: Word Document (DOCX)]\n\n"
            text += "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])
            
        elif filename_lower.endswith((".xlsx", ".xls")):
            wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
            text += f"[DOCUMENT TYPE: Excel Spreadsheet ({len(wb.sheetnames)} Sheets)]\n\n"
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                text += f"--- Sheet: {sheet} ---\n"
                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) for cell in row if cell is not None])
                    if row_text:
                        text += row_text + "\n"
                text += "\n"
                
        elif filename_lower.endswith((".pptx", ".ppt")):
            prs = Presentation(io.BytesIO(content))
            text += f"[DOCUMENT TYPE: PowerPoint Presentation ({len(prs.slides)} Slides)]\n\n"
            for i, slide in enumerate(prs.slides):
                text += f"--- Slide {i+1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
                text += "\n"
                
        elif filename_lower.endswith((".txt", ".md", ".csv")):
            text += f"[DOCUMENT TYPE: Text / Code File ({filename})]\n\n"
            try:
                text += content.decode("utf-8")
            except UnicodeDecodeError:
                text += content.decode("latin-1", errors="replace")
            
        elif filename_lower.endswith((".png", ".jpg", ".jpeg", ".webp")):
            try:
                mime_type = f"image/{filename_lower.split('.')[-1]}"
                if mime_type == "image/jpg":
                    mime_type = "image/jpeg"
                
                # Preprocess image bytes with OpenCV to enhance contrast/deskew
                processed_bytes = content
                try:
                    processed_bytes = preprocess_for_ocr(content)
                except Exception as prep_err:
                    print(f"OpenCV preprocessing error (continuing with raw bytes): {prep_err}")
                
                vision_analysis = extract_text_with_vision_api(processed_bytes, mime_type=mime_type)
                text += vision_analysis
            except Exception as e:
                print(f"Image Vision failed for {filename}: {e}")
                text += f"[ATTACHED IMAGE FILE: {filename}]\nUploaded image document ready for analysis."
        else:
            text += f"Attached file: {filename} (Content ready for analysis)"
            
    except Exception as e:
        print(f"Text extraction error for {filename}: {e}")
        text += f"Attached document: {filename} (Uploaded and processed)."
        
    return text

def chunk_text(text: str) -> list[str]:
    """
    Split text into smaller chunks for vectorization using LangChain's splitter.
    """
    if not text.strip():
        return []
        
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len,
        is_separator_regex=False,
    )
    
    chunks = text_splitter.split_text(text)
    if not chunks and text.strip():
        chunks = [text.strip()]
    return chunks

async def process_upload(file: UploadFile, user_id: str, chat_id: str = None):
    """
    Process an uploaded file: extract text, chunk it, embed it, and store it.
    """
    # 1. Extract text
    text = await extract_text(file)
    
    file.file.seek(0, 2)
    file_size = file.file.tell()
    file.file.seek(0)
    
    if not text.strip():
        text = f"--- UPLOADED FILE: {file.filename} ---\nAttached document available for analysis."

    print("\n========== EXTRACTED TEXT ==========")
    print(text[:500])
    print("===================================\n")
        
    # 2. Chunk text
    chunks = chunk_text(text)
    if not chunks:
        chunks = [text]

    print("\n========== CHUNKS ==========")
    print(f"Number of chunks: {len(chunks)}")
    if chunks:
        print(f"Sample chunk 1:\n{chunks[0][:300]}")
    print("============================\n")
        
    # 3. Embed chunks
    embeddings = get_embeddings(chunks)
    print(f"[Embedding Generation] Generated {len(embeddings)} embeddings for {len(chunks)} chunks.")
    
    # 4. Store in ChromaDB
    document_id = str(uuid.uuid4())
    success = store_chunks(user_id, document_id, file.filename, chunks, embeddings, chat_id=chat_id)
    print(f"[ChromaDB Storage] Store status: {success} for user_id='{user_id}', chat_id='{chat_id}'")
    
    if success:
        return {
            "filename": file.filename,
            "document_id": document_id,
            "chunks_stored": len(chunks),
            "size": file_size,
            "status": "success"
        }
    return {"filename": file.filename, "size": file_size, "status": "failed", "reason": "Database storage failed"}
